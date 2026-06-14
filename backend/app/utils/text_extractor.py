import io
import json
import base64
import PyPDF2
import docx
import openpyxl
from PIL import Image


def extract_text_from_pdf(file_bytes: bytes) -> str:
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = docx.Document(io.BytesIO(file_bytes))
    full_text = []
    
    # Extract paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
            
    # Extract tables
    for table in doc.tables:
        full_text.append("")  # Blank line separator
        for row in table.rows:
            # Join cells with pipes, replacing newlines within cells
            row_text = " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
            if row_text.strip().replace("|", "").strip():
                full_text.append(row_text)
                
    return "\n".join(full_text)


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def _gemini_vision(file_bytes: bytes, filename: str, prompt: str, max_tokens: int = 4096) -> str:
    """Send an image to Gemini vision model with a given prompt."""
    import google.generativeai as genai
    from app.config import settings

    genai.configure(api_key=settings.gemini_api_key)
    model = genai.GenerativeModel("gemini-2.5-flash")
    
    img = Image.open(io.BytesIO(file_bytes))
    
    generation_config = genai.types.GenerationConfig(
        max_output_tokens=max_tokens,
    )
    
    response = model.generate_content([prompt, img], generation_config=generation_config)
    return response.text


def extract_text_from_image(file_bytes: bytes, filename: str) -> str:
    return _gemini_vision(
        file_bytes, filename,
        "Extract and return ALL text visible in this image. Include all table content, headings, labels, room numbers, and any structured data as plain text.",
    )


def extract_timetable_from_image(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract structured timetable entries directly from an image using Groq vision."""
    import json as _json
    prompt = (
        "This is a university timetable image. "
        "Extract every class/subject slot and return ONLY a valid JSON array, no explanation. "
        "Each object must have exactly these keys: "
        "day (Monday/Tuesday/Wednesday/Thursday/Friday/Saturday/Sunday), "
        "subject (course code and name), time (start time e.g. 09:00 AM), "
        "end_time (e.g. 10:00 AM), room (room/lab number if visible, else empty string), "
        "instructor (teacher name/code if visible, else empty string). "
        "If a cell spans multiple days, repeat the entry for each day. "
        "Return only the JSON array."
    )
    raw = _gemini_vision(file_bytes, filename, prompt, max_tokens=4096)
    try:
        start = raw.find("[")
        end = raw.rfind("]") + 1
        return _json.loads(raw[start:end])
    except Exception:
        return []


def extract_timetable_from_pdf(file_bytes: bytes) -> list[dict]:
    """Convert first PDF page to image then use vision extraction."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        page = doc[0]
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("jpeg")
        # Use a temp filename to signal jpeg mime
        return extract_timetable_from_image(img_bytes, "timetable.jpg")
    except ImportError:
        return []  # PyMuPDF not installed, fall back to text extraction


def extract_text_from_pptx(file_bytes: bytes) -> str:
    import pptx
    prs = pptx.Presentation(io.BytesIO(file_bytes))
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"[Slide {i+1}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)


def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext in ("docx", "doc"):
        return extract_text_from_docx(file_bytes)
    elif ext in ("xlsx", "xls"):
        return extract_text_from_xlsx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return extract_text_from_pptx(file_bytes)
    elif ext in ("jpg", "jpeg", "png"):
        return extract_text_from_image(file_bytes, filename)
    elif ext == "json":
        data = json.loads(file_bytes.decode("utf-8"))
        return json.dumps(data, indent=2)
    else:
        return file_bytes.decode("utf-8", errors="ignore")


def parse_timetable_json(file_bytes: bytes) -> dict:
    return json.loads(file_bytes.decode("utf-8"))
